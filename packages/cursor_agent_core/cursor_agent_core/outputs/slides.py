"""Generate PowerPoint slide decks from structured content."""

from __future__ import annotations

import json
from typing import Any

from cursor_sdk import CustomTool, CustomToolContext
from pptx import Presentation
from pptx.util import Inches, Pt

from cursor_agent_core.paths.file_helpers import (
    relative_to_project,
    resolve_under_root,
    slugify,
    unique_path,
)
from cursor_agent_core.paths.project_context import get_output_directory, get_project_root

_VALID_LAYOUTS = {"title", "bullets", "two_column", "image"}


def _parse_slides(raw_slides: Any) -> list[dict[str, Any]]:
    if not isinstance(raw_slides, list) or not raw_slides:
        raise ValueError("slides is required (non-empty list)")

    slides: list[dict[str, Any]] = []
    for item in raw_slides:
        if not isinstance(item, dict):
            raise ValueError("each slide must be an object")
        layout = str(item.get("layout") or "bullets").strip().lower()
        if layout not in _VALID_LAYOUTS:
            raise ValueError(f"unsupported slide layout: {layout}")
        slides.append({**item, "layout": layout})
    return slides


def _set_title(slide: Any, text: str) -> None:
    if slide.shapes.title is not None:
        slide.shapes.title.text = text


def _add_title_slide(prs: Presentation, slide_data: dict[str, Any]) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    _set_title(slide, str(slide_data.get("title") or ""))
    subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    if subtitle is not None:
        subtitle.text = str(slide_data.get("subtitle") or "")


def _add_bullets_slide(prs: Presentation, slide_data: dict[str, Any]) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    _set_title(slide, str(slide_data.get("title") or ""))

    body = slide.placeholders[1]
    text_frame = body.text_frame
    text_frame.clear()
    bullets = slide_data.get("bullets") or []
    if not isinstance(bullets, list):
        raise ValueError("bullets layout requires a bullets array")

    for index, bullet in enumerate(bullets):
        text = str(bullet).strip()
        if not text:
            continue
        if index == 0:
            text_frame.text = text
        else:
            text_frame.add_paragraph().text = text


def _add_two_column_slide(prs: Presentation, slide_data: dict[str, Any]) -> None:
    layout = prs.slide_layouts[3] if len(prs.slide_layouts) > 3 else prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    _set_title(slide, str(slide_data.get("title") or ""))

    left_items = slide_data.get("left_bullets") or []
    right_items = slide_data.get("right_bullets") or []
    if not isinstance(left_items, list) or not isinstance(right_items, list):
        raise ValueError("two_column layout requires left_bullets and right_bullets arrays")

    if len(slide.placeholders) > 1:
        left = slide.placeholders[1].text_frame
        left.clear()
        for index, item in enumerate(left_items):
            text = str(item).strip()
            if not text:
                continue
            if index == 0:
                left.text = text
            else:
                left.add_paragraph().text = text

    if len(slide.placeholders) > 2:
        right = slide.placeholders[2].text_frame
        right.clear()
        for index, item in enumerate(right_items):
            text = str(item).strip()
            if not text:
                continue
            if index == 0:
                right.text = text
            else:
                right.add_paragraph().text = text


def _add_image_slide(prs: Presentation, slide_data: dict[str, Any]) -> None:
    layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    _set_title(slide, str(slide_data.get("title") or ""))

    image_path_raw = str(slide_data.get("image_path") or "").strip()
    if not image_path_raw:
        raise ValueError("image layout requires image_path")

    image_path = resolve_under_root(image_path_raw, get_project_root())
    if not image_path.exists():
        raise ValueError(f"image not found: {image_path_raw}")

    slide.shapes.add_picture(
        str(image_path),
        Inches(1.0),
        Inches(1.5),
        width=Inches(8.0),
    )

    caption = str(slide_data.get("caption") or "").strip()
    if caption:
        box = slide.shapes.add_textbox(Inches(1.0), Inches(6.5), Inches(8.0), Inches(0.8))
        paragraph = box.text_frame.paragraphs[0]
        paragraph.text = caption
        paragraph.font.size = Pt(14)


def _build_deck(title: str, slides: list[dict[str, Any]]) -> Presentation:
    prs = Presentation()
    for slide_data in slides:
        layout = slide_data["layout"]
        if layout == "title":
            _add_title_slide(prs, slide_data)
        elif layout == "bullets":
            _add_bullets_slide(prs, slide_data)
        elif layout == "two_column":
            _add_two_column_slide(prs, slide_data)
        elif layout == "image":
            _add_image_slide(prs, slide_data)
    return prs


def create_slide_deck(
    args: dict[str, Any],
    _ctx: CustomToolContext,
) -> str:
    """Create a PowerPoint deck in the session Slides/ folder."""
    title = str(args.get("title") or "").strip()
    if not title:
        raise ValueError("title is required")

    slides = _parse_slides(args.get("slides"))
    directory = get_output_directory('slides')
    deck_path = unique_path(directory, f"{slugify(title)}.pptx")

    presentation = _build_deck(title, slides)
    presentation.save(deck_path)

    payload = {
        "title": title,
        "slides_directory": relative_to_project(directory, get_project_root()),
        "pptx_path": relative_to_project(deck_path, get_project_root()),
        "slide_count": len(slides),
    }
    return json.dumps(payload, indent=2)


CREATE_SLIDE_DECK_TOOL = CustomTool(
    execute=create_slide_deck,
    description=(
        "Create a PowerPoint slide deck in the session Slides/ folder. Supports title, bullets, "
        "two_column, and image layouts. Image slides can reference session Diagrams/ PNGs."
    ),
    input_schema={
        "type": "object",
        "properties": {
            "title": {"type": "string", "description": "Deck title."},
            "slides": {
                "type": "array",
                "description": "Ordered slide definitions.",
                "items": {
                    "type": "object",
                    "properties": {
                        "layout": {
                            "type": "string",
                            "enum": ["title", "bullets", "two_column", "image"],
                        },
                        "title": {"type": "string"},
                        "subtitle": {"type": "string"},
                        "bullets": {
                            "type": "array",
                            "items": {"type": "string"},
                        },
                        "left_bullets": {
                            "type": "array",
                            "items": {"type": "string"},
                        },
                        "right_bullets": {
                            "type": "array",
                            "items": {"type": "string"},
                        },
                        "image_path": {
                            "type": "string",
                            "description": "Project-relative image path.",
                        },
                        "caption": {"type": "string"},
                    },
                    "required": ["layout"],
                },
            },
        },
        "required": ["title", "slides"],
    },
)
