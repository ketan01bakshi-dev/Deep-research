"""PDF viewer and paper listing for research sessions."""

from __future__ import annotations

import streamlit as st

from deep_research.session_context import ResearchSession, list_session_files
from deep_research.web.components import render_empty_state


def render_papers_tab(session: ResearchSession) -> None:
    papers = [p for p in list_session_files(session, "Papers") if "summaries" not in p.parts]
    if papers:
        for pdf_path in papers:
            st.markdown(f"**{pdf_path.name}**")
            try:
                st.pdf(str(pdf_path))
            except Exception:
                st.info("Inline PDF preview unavailable for this file.")
            with open(pdf_path, "rb") as handle:
                st.download_button(
                    label=f"Download {pdf_path.name}",
                    data=handle.read(),
                    file_name=pdf_path.name,
                    mime="application/pdf",
                    key=f"dl_{session.id}_{pdf_path.name}",
                )
            st.divider()
    else:
        render_empty_state("No PDFs downloaded for this session yet.")

    summaries_dir = session.root / "Papers" / "summaries"
    if summaries_dir.exists():
        summary_files = sorted(summaries_dir.glob("*.md"))
        if summary_files:
            st.markdown("### Paper summaries")
            for summary_path in summary_files:
                st.markdown(f"**{summary_path.name}**")
                st.markdown(summary_path.read_text(encoding="utf-8")[:8000])
                st.divider()

    _render_rag_ask(session)


def _render_rag_ask(session: ResearchSession) -> None:
    """Ask this research — grounded query over session corpus."""
    st.markdown("### Ask this research")
    question = st.text_input(
        "Question grounded in session PDFs/reports",
        key=f"rag_q_{session.id}",
        placeholder="What did the papers say about X?",
    )
    col1, col2, col3 = st.columns(3)
    with col1:
        top_k = st.slider("Top K", 1, 10, 5, key=f"rag_topk_{session.id}")
    with col2:
        rebuild = st.checkbox("Rebuild index", value=False, key=f"rag_rebuild_{session.id}")
    with col3:
        run_query = st.button("Query corpus", key=f"rag_btn_{session.id}")
    if run_query:
        if not question.strip():
            st.warning("Enter a question.")
            return
        from deep_research.rag import build_session_corpus, format_corpus_context, query_session_corpus

        if rebuild:
            build_session_corpus(session.id)
        hits = query_session_corpus(session.id, question.strip(), top_k=top_k)
        if not hits:
            st.info("No matching excerpts in session corpus.")
        else:
            st.caption(f"Retrieved {len(hits)} excerpt(s)")
            for index, hit in enumerate(hits, start=1):
                with st.expander(f"Hit {index}: {hit.get('source', 'unknown')}"):
                    st.markdown(str(hit.get("text") or "")[:2000])
            st.markdown(format_corpus_context(hits))


def render_report_tab(session: ResearchSession) -> None:
    reports = list_session_files(session, "Reports")
    if not reports:
        render_empty_state("No report files for this session.")
        return

    for report_path in reports:
        st.markdown(f"**{report_path.name}**")
        if report_path.suffix.lower() == ".md":
            st.markdown(report_path.read_text(encoding="utf-8"))
        from deep_research.web.export import report_pdf_for_session

        pdf_bytes = report_pdf_for_session(session)
        if pdf_bytes and report_path.suffix.lower() == ".md":
            st.download_button(
                "Download report as PDF",
                data=pdf_bytes,
                file_name=f"{report_path.stem}.pdf",
                mime="application/pdf",
                key=f"dl_report_pdf_{session.id}_{report_path.name}",
            )
        with open(report_path, "rb") as handle:
            mime = (
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                if report_path.suffix.lower() == ".docx"
                else "text/markdown"
            )
            st.download_button(
                label=f"Download {report_path.name}",
                data=handle.read(),
                file_name=report_path.name,
                mime=mime,
                key=f"dl_report_{session.id}_{report_path.name}",
            )
        st.divider()


def render_mindmap_tab(session: ResearchSession) -> None:
    diagrams = [
        p
        for p in list_session_files(session, "Diagrams")
        if p.suffix.lower() in {".png", ".svg", ".jpg", ".jpeg", ".webp"}
    ]
    if not diagrams:
        render_empty_state("No mind map or diagram images for this session.")
        return

    for image_path in diagrams:
        st.markdown(f"**{image_path.name}**")
        st.image(str(image_path), use_container_width=True)
        with open(image_path, "rb") as handle:
            st.download_button(
                label=f"Download {image_path.name}",
                data=handle.read(),
                file_name=image_path.name,
                key=f"dl_diagram_{session.id}_{image_path.name}",
            )
        st.divider()


def render_slides_tab(session: ResearchSession) -> None:
    slides = [p for p in list_session_files(session, "Slides") if p.suffix.lower() == ".pptx"]
    if not slides:
        render_empty_state("No slide decks for this session.")
        return

    for slide_path in slides:
        st.markdown(f"**{slide_path.name}**")
        try:
            from pptx import Presentation

            prs = Presentation(str(slide_path))
            st.caption(f"{len(prs.slides)} slides")
            for index, slide in enumerate(prs.slides, start=1):
                title = ""
                if slide.shapes.title and slide.shapes.title.text:
                    title = slide.shapes.title.text.strip()
                label = title or f"Slide {index}"
                st.markdown(f"- {label}")
        except Exception:
            st.caption("Slide outline unavailable — download to view.")
        with open(slide_path, "rb") as handle:
            st.download_button(
                label=f"Download {slide_path.name}",
                data=handle.read(),
                file_name=slide_path.name,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key=f"dl_slides_{session.id}_{slide_path.name}",
            )
        st.divider()
